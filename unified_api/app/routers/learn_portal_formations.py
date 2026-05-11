from fastapi import APIRouter, HTTPException, Request, Response, UploadFile, File
from pydantic import BaseModel
from typing import Optional, Any
from psycopg.rows import dict_row
from io import BytesIO
import uuid
import json
import re
import unicodedata
import html
import os
from difflib import SequenceMatcher

from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm
from reportlab.platypus import Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import ParagraphStyle

from app.routers.skills_portal_common import get_conn
from app.routers.learn_portal_common import learn_require_user, learn_fetch_profile
from app.routers.skills_portal_pdf_common import build_pdf_document, build_pdf_styles

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation as PptxPresentation
except Exception:
    PptxPresentation = None

router = APIRouter()


# ======================================================
# Helpers Learn
# ======================================================

def _learn_require_profile(cur, u: dict, id_effectif: str) -> dict:
    eid = (id_effectif or "").strip()
    if not eid:
        raise HTTPException(status_code=400, detail="id_effectif manquant.")

    profile = learn_fetch_profile(
        cur,
        id_effectif=eid,
        email=(u.get("email") or ""),
        is_super_admin=bool(u.get("is_super_admin")),
    )

    oid = (profile.get("id_owner") or "").strip()
    if not oid:
        raise HTTPException(status_code=403, detail="Profil Learn sans owner.")

    role = (profile.get("role_code") or "user").strip().lower()
    if role not in ("admin", "supervisor", "user"):
        role = "user"

    profile["role_code"] = role
    return profile


def _role_rank(role_code: str) -> int:
    c = (role_code or "").strip().lower()
    if c == "admin":
        return 3
    if c == "supervisor":
        return 2
    return 1


def _learn_require_min_role(profile: dict, min_role: str) -> None:
    if _role_rank(profile.get("role_code")) < _role_rank(min_role):
        raise HTTPException(status_code=403, detail="Accès refusé : droits insuffisants.")


def _json_list(value: Any) -> list:
    if value is None:
        return []

    if isinstance(value, list):
        return [str(x).strip() for x in value if str(x or "").strip()]

    if isinstance(value, tuple):
        return [str(x).strip() for x in value if str(x or "").strip()]

    if isinstance(value, str):
        s = value.strip()
        if not s:
            return []

        try:
            js = json.loads(s)
            if isinstance(js, list):
                return [str(x).strip() for x in js if str(x or "").strip()]
        except Exception:
            return []

    return []


def _jsonb_param(value: Any) -> str:
    arr = _json_list(value)
    return json.dumps(arr, ensure_ascii=False)


def _clean_text(value: Any, max_len: int = 20000) -> str:
    txt = str(value or "").replace("\x00", " ").strip()
    txt = re.sub(r"\r\n?", "\n", txt)
    txt = re.sub(r"[ \t]+", " ", txt)
    if len(txt) > max_len:
        txt = txt[:max_len].rsplit(" ", 1)[0].strip()
    return txt


def _safe_float(value: Any) -> Optional[float]:
    if value is None:
        return None

    s = str(value).replace(",", ".").strip()
    if not s:
        return None

    try:
        return float(s)
    except Exception:
        return None

def _normalize_type_formation(value: Any) -> str:
    v = str(value or "").strip().lower()

    if v == "certifiante":
        return "Certifiante"

    if v in ("diplomante", "diplômante"):
        return "Diplomante"

    if v in ("non certifiante", "non-certifiante", "non certifiant", "non-certifiant"):
        return "Non Certifiante"

    return "Non Certifiante"

def _safe_int(value: Any) -> Optional[int]:
    if value is None:
        return None

    s = str(value).strip()
    if not s:
        return None

    try:
        return int(float(s.replace(",", ".")))
    except Exception:
        return None

def _duration_text(value: Any) -> Optional[str]:
    n = _safe_float(value)

    if n is None:
        return None

    if float(n).is_integer():
        return str(int(n))

    return str(n).replace(",", ".").strip()


def _next_plan_code(cur, oid: str) -> str:
    cur.execute("SELECT pg_advisory_xact_lock(hashtext(%s))", (f"learn_plan_code:{oid}:PL",))

    cur.execute(
        """
        SELECT COALESCE(
          MAX((regexp_match(codification, '^PL([0-9]{5})$'))[1]::int),
          0
        ) AS max_n
        FROM public.tbl_plan_pedagogique
        WHERE id_owner = %s
          AND codification ~ '^PL[0-9]{5}$'
        """,
        (oid,),
    )

    r = cur.fetchone() or {}
    max_n_raw = r.get("max_n")
    max_n = int(max_n_raw) if max_n_raw is not None else 0
    nxt = max_n + 1

    if nxt > 99999:
        raise HTTPException(status_code=400, detail="Limite de numérotation atteinte (PL99999).")

    return f"PL{nxt:05d}"

def _next_form_code(cur, oid: str) -> str:
    cur.execute("SELECT pg_advisory_xact_lock(hashtext(%s))", (f"learn_form_code:{oid}:FC",))

    cur.execute(
        """
        SELECT COALESCE(
          MAX((regexp_match(code, '^FC([0-9]{5})$'))[1]::int),
          0
        ) AS max_n
        FROM public.tbl_fiche_formation
        WHERE id_owner = %s
          AND code ~ '^FC[0-9]{5}$'
        """,
        (oid,),
    )

    r = cur.fetchone() or {}
    max_n_raw = r.get("max_n")
    max_n = int(max_n_raw) if max_n_raw is not None else 0
    nxt = max_n + 1

    if nxt > 99999:
        raise HTTPException(status_code=400, detail="Limite de numérotation atteinte (FC99999).")

    return f"FC{nxt:05d}"


def _formation_exists_owner(cur, oid: str, id_form: str) -> bool:
    cur.execute(
        """
        SELECT 1
        FROM public.tbl_fiche_formation
        WHERE id_owner = %s
          AND id_form = %s
        LIMIT 1
        """,
        (oid, id_form),
    )
    return cur.fetchone() is not None


def _fetch_owner_logo_bytes(cur, oid: str) -> Optional[bytes]:
    cur.execute(
        """
        SELECT logo_bytes
        FROM public.tbl_studio_owner_logo
        WHERE id_owner = %s
          AND COALESCE(archive, FALSE) = FALSE
        ORDER BY date_maj DESC, date_creation DESC
        LIMIT 1
        """,
        (oid,),
    )

    row = cur.fetchone() or {}
    raw = row.get("logo_bytes")

    if raw is None:
        return None

    try:
        return bytes(raw)
    except Exception:
        return raw


def _pdf_safe_filename_part(v: Any, max_len: int = 120) -> str:
    s = str(v or "").strip()
    if not s:
        return "document"

    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    s = re.sub(r"[^A-Za-z0-9._ -]+", "", s)
    s = re.sub(r"\s+", " ", s).strip(" ._-")

    if not s:
        s = "document"

    if len(s) > max_len:
        s = s[:max_len].strip(" ._-")

    return s or "document"


def _pdf_latin1_safe(v: Any) -> str:
    return str(v or "").encode("latin-1", "replace").decode("latin-1")


def _resolve_labels(cur, table: str, id_col: str, ids: list) -> list:
    clean_ids = [str(x).strip() for x in ids if str(x or "").strip()]
    if not clean_ids:
        return []

    if table not in ("tbl_mod_form", "tbl_met_peda", "tbl_met_eval"):
        return []

    cur.execute(
        f"""
        SELECT
          {id_col} AS id,
          titre,
          titre_court,
          description,
          ordre_affichage
        FROM public.{table}
        WHERE {id_col} = ANY(%s)
          AND COALESCE(masque, FALSE) = FALSE
        ORDER BY COALESCE(ordre_affichage, 999999), lower(COALESCE(titre, titre_court, {id_col}))
        """,
        (clean_ids,),
    )

    rows = cur.fetchall() or []
    by_id = {str(r.get("id") or "").strip(): r for r in rows}

    return [
        by_id[x]
        for x in clean_ids
        if x in by_id
    ]


def _resolve_competences(cur, oid: str, ids: list) -> list:
    clean_ids = [str(x).strip() for x in ids if str(x or "").strip()]
    if not clean_ids:
        return []

    cur.execute(
        """
        SELECT
          c.id_comp,
          c.code,
          c.intitule,
          c.domaine,
          dc.titre_court AS domaine_titre_court,
          dc.titre AS domaine_titre,
          dc.couleur AS domaine_couleur,
          c.etat
        FROM public.tbl_competence c
        LEFT JOIN public.tbl_domaine_competence dc
          ON dc.id_domaine_competence = c.domaine
         AND COALESCE(dc.masque, FALSE) = FALSE
        WHERE c.id_owner = %s
          AND c.id_comp = ANY(%s)
          AND COALESCE(c.masque, FALSE) = FALSE
        ORDER BY lower(COALESCE(c.code, '')), lower(COALESCE(c.intitule, ''))
        """,
        (oid, clean_ids),
    )

    rows = cur.fetchall() or []
    by_id = {str(r.get("id_comp") or "").strip(): r for r in rows}

    return [
        by_id[x]
        for x in clean_ids
        if x in by_id
    ]

def _normalize_competence_ids(value: Any) -> list:
    ids = _json_list(value)
    out = []
    seen = set()

    for raw in ids:
        cid = str(raw or "").strip()
        if not cid or cid in seen:
            continue
        seen.add(cid)
        out.append(cid)

    return out


def _ensure_competences_owner(cur, oid: str, ids: list) -> list:
    clean_ids = _normalize_competence_ids(ids)
    if not clean_ids:
        return []

    cur.execute(
        """
        SELECT id_comp
        FROM public.tbl_competence
        WHERE id_owner = %s
          AND id_comp = ANY(%s)
          AND COALESCE(masque, FALSE) = FALSE
        """,
        (oid, clean_ids),
    )

    valid = {str(r.get("id_comp") or "").strip() for r in (cur.fetchall() or [])}
    return [cid for cid in clean_ids if cid in valid]


def _fetch_contenu_row(cur, oid: str, id_form: str, id_ligne_contenu: str) -> dict:
    cur.execute(
        """
        SELECT
          l.id_ligne_contenu,
          l.titre_sequence,
          l.objectif,
          l.contenu,
          l.id_competence,
          l.competences_liees,
          l.position
        FROM public.tbl_contenu_ligne l
        WHERE l.id_owner = %s
          AND l.id_form = %s
          AND l.id_ligne_contenu = %s
          AND COALESCE(l.archive, FALSE) = FALSE
        LIMIT 1
        """,
        (oid, id_form, id_ligne_contenu),
    )

    row = cur.fetchone()
    if not row:
        raise HTTPException(status_code=404, detail="Contenu introuvable.")

    comp_ids = _normalize_competence_ids(row.get("competences_liees"))
    if not comp_ids and row.get("id_competence"):
        comp_ids = [str(row.get("id_competence")).strip()]

    row["competences_liees_ids"] = comp_ids
    row["competences_liees_items"] = _resolve_competences(cur, oid, comp_ids)

    return row


def _renumber_contenus(cur, oid: str, id_form: str) -> None:
    cur.execute(
        """
        SELECT id_ligne_contenu
        FROM public.tbl_contenu_ligne
        WHERE id_owner = %s
          AND id_form = %s
          AND COALESCE(archive, FALSE) = FALSE
        ORDER BY COALESCE(position, 999999), titre_sequence
        """,
        (oid, id_form),
    )

    rows = cur.fetchall() or []

    for idx, r in enumerate(rows, start=1):
        cur.execute(
            """
            UPDATE public.tbl_contenu_ligne
            SET position = %s,
                date_modification = NOW()
            WHERE id_owner = %s
              AND id_form = %s
              AND id_ligne_contenu = %s
            """,
            (idx, oid, id_form, r.get("id_ligne_contenu")),
        )

def _ensure_contenus_owner(cur, oid: str, id_form: str, ids: list) -> list:
    clean_ids = [str(x or "").strip() for x in (ids or []) if str(x or "").strip()]
    if not clean_ids:
        return []

    cur.execute(
        """
        SELECT id_ligne_contenu
        FROM public.tbl_contenu_ligne
        WHERE id_owner = %s
          AND id_form = %s
          AND id_ligne_contenu = ANY(%s)
          AND COALESCE(archive, FALSE) = FALSE
        """,
        (oid, id_form, clean_ids),
    )

    valid = {str(r.get("id_ligne_contenu") or "").strip() for r in (cur.fetchall() or [])}
    return [cid for cid in clean_ids if cid in valid]


def _fetch_plan_detail(cur, oid: str, id_form: str, id_plan_peda: str) -> dict:
    cur.execute(
        """
        SELECT
          p.id_plan_peda,
          p.codification,
          p.titre,
          p.commentaire,
          p.modalite_generale,
          p.chemin_sharepoint,
          p.date_creation,
          COALESCE(
            SUM(
              CASE
                WHEN trim(COALESCE(b.duree::text, '')) ~ '^[0-9]+([\\.,][0-9]+)?$'
                THEN replace(trim(b.duree::text), ',', '.')::numeric
                ELSE 0
              END
            ),
            0
          ) AS duree_totale,
          COUNT(DISTINCT b.id_bloc_peda) AS nb_blocs
        FROM public.tbl_plan_pedagogique p
        LEFT JOIN public.tbl_bloc_pedagogique b
          ON b.id_owner = p.id_owner
         AND b.id_plan_peda = p.id_plan_peda
         AND COALESCE(b.archive, FALSE) = FALSE
        WHERE p.id_owner = %s
          AND p.id_form = %s
          AND p.id_plan_peda = %s
          AND COALESCE(p.archive, FALSE) = FALSE
        GROUP BY
          p.id_plan_peda,
          p.codification,
          p.titre,
          p.commentaire,
          p.modalite_generale,
          p.chemin_sharepoint,
          p.date_creation
        LIMIT 1
        """,
        (oid, id_form, id_plan_peda),
    )

    plan = cur.fetchone()
    if not plan:
        raise HTTPException(status_code=404, detail="Plan pédagogique introuvable.")

    cur.execute(
        """
        SELECT
          b.id_bloc_peda,
          b.titre,
          b.objectif,
          b.duree,
          b.modalite_intervention,
          b.observations,
          b.position
        FROM public.tbl_bloc_pedagogique b
        WHERE b.id_owner = %s
          AND b.id_form = %s
          AND b.id_plan_peda = %s
          AND COALESCE(b.archive, FALSE) = FALSE
        ORDER BY COALESCE(b.position, 999999), lower(COALESCE(b.titre, ''))
        """,
        (oid, id_form, id_plan_peda),
    )

    blocs = cur.fetchall() or []

    for b in blocs:
        cur.execute(
            """
            SELECT
              s.id_sequence_bloc,
              s.id_ligne_contenu,
              s.position,
              l.titre_sequence,
              l.objectif,
              l.contenu,
              l.id_competence,
              l.competences_liees
            FROM public.tbl_sequence_bloc_pedagogique s
            JOIN public.tbl_contenu_ligne l
              ON l.id_owner = s.id_owner
             AND l.id_ligne_contenu = s.id_ligne_contenu
             AND COALESCE(l.archive, FALSE) = FALSE
            WHERE s.id_owner = %s
              AND s.id_bloc_peda = %s
              AND COALESCE(s.archive, FALSE) = FALSE
            ORDER BY COALESCE(s.position, 999999), lower(COALESCE(l.titre_sequence, ''))
            """,
            (oid, b.get("id_bloc_peda")),
        )

        seqs = cur.fetchall() or []

        for s in seqs:
            comp_ids = _normalize_competence_ids(s.get("competences_liees"))
            if not comp_ids and s.get("id_competence"):
                comp_ids = [str(s.get("id_competence")).strip()]

            s["competences_liees_ids"] = comp_ids
            s["competences_liees_items"] = _resolve_competences(cur, oid, comp_ids)

        b["sequences"] = seqs

    plan["blocs"] = blocs
    return plan


def _archive_plan_blocs_and_sequences(cur, oid: str, id_form: str, id_plan_peda: str) -> None:
    cur.execute(
        """
        UPDATE public.tbl_sequence_bloc_pedagogique s
        SET archive = TRUE,
            date_modification = NOW()
        FROM public.tbl_bloc_pedagogique b
        WHERE b.id_owner = s.id_owner
          AND b.id_bloc_peda = s.id_bloc_peda
          AND b.id_owner = %s
          AND b.id_form = %s
          AND b.id_plan_peda = %s
          AND COALESCE(s.archive, FALSE) = FALSE
        """,
        (oid, id_form, id_plan_peda),
    )

    cur.execute(
        """
        UPDATE public.tbl_bloc_pedagogique
        SET archive = TRUE,
            date_modification = NOW()
        WHERE id_owner = %s
          AND id_form = %s
          AND id_plan_peda = %s
          AND COALESCE(archive, FALSE) = FALSE
        """,
        (oid, id_form, id_plan_peda),
    )


def _insert_plan_blocs(cur, oid: str, id_form: str, id_plan_peda: str, blocs: Optional[list]) -> None:
    rows = blocs or []

    for idx, raw in enumerate(rows, start=1):
        if raw is None:
            continue

        item = raw if isinstance(raw, dict) else raw.dict()

        titre = _clean_text(item.get("titre") or f"Séquence {idx}", 500)
        duree = _duration_text(item.get("duree"))
        modalite = _clean_text(item.get("modalite_intervention"), 250)
        objectif = _clean_text(item.get("objectif"), 1200)
        observations = _clean_text(item.get("observations"), 2000)

        contenu_ids = _ensure_contenus_owner(cur, oid, id_form, item.get("contenus") or [])

        bid = str(uuid.uuid4())

        cur.execute(
            """
            INSERT INTO public.tbl_bloc_pedagogique
              (
                id_bloc_peda,
                id_owner,
                id_form,
                id_plan_peda,
                titre,
                objectif,
                duree,
                modalite_intervention,
                observations,
                position,
                archive,
                date_creation,
                date_modification
              )
            VALUES
              (
                %s, %s, %s, %s,
                %s, %s, %s, %s, %s,
                %s,
                FALSE,
                NOW(),
                NOW()
              )
            """,
            (
                bid,
                oid,
                id_form,
                id_plan_peda,
                titre,
                objectif,
                duree,
                modalite,
                observations,
                idx,
            ),
        )

        for seq_idx, id_ligne in enumerate(contenu_ids, start=1):
            cur.execute(
                """
                INSERT INTO public.tbl_sequence_bloc_pedagogique
                  (
                    id_sequence_bloc,
                    id_owner,
                    id_bloc_peda,
                    id_ligne_contenu,
                    position,
                    archive,
                    date_creation,
                    date_modification
                  )
                VALUES
                  (
                    %s, %s, %s, %s,
                    %s,
                    FALSE,
                    NOW(),
                    NOW()
                  )
                """,
                (
                    str(uuid.uuid4()),
                    oid,
                    bid,
                    id_ligne,
                    seq_idx,
                ),
            )

def _sync_formation_prerequis(cur, oid: str, id_form: str, prerequis: Optional[list]) -> None:
    """
    Synchronise les prérequis évaluables d'une fiche formation.

    Principe Novoskill :
    - pas de suppression physique ;
    - archivage logique des anciennes lignes ;
    - réactivation/update si id_prerequis fourni ;
    - création si nouveau prérequis.
    """
    items = prerequis or []

    cur.execute(
        """
        UPDATE public.tbl_fiche_formation_prerequis
        SET archive = TRUE,
            date_modification = NOW()
        WHERE id_owner = %s
          AND id_form = %s
          AND COALESCE(archive, FALSE) = FALSE
        """,
        (oid, id_form),
    )

    ordre = 1

    for raw in items:
        if raw is None:
            continue

        if isinstance(raw, dict):
            item = raw
        else:
            item = raw.dict()

        titre = _clean_text(item.get("titre"), 800)
        if not titre:
            continue

        r1 = _clean_text(item.get("r1") or "Je ne maîtrise pas", 300)
        r2 = _clean_text(item.get("r2") or "J’ai besoin d’assistance", 300)
        r3 = _clean_text(item.get("r3") or "", 300)

        try:
            item_ordre = int(item.get("ordre_affichage") or ordre)
        except Exception:
            item_ordre = ordre

        pid = (item.get("id_prerequis") or "").strip()

        if pid:
            cur.execute(
                """
                UPDATE public.tbl_fiche_formation_prerequis
                SET titre = %s,
                    r1 = %s,
                    r2 = %s,
                    r3 = %s,
                    ordre_affichage = %s,
                    archive = FALSE,
                    date_modification = NOW()
                WHERE id_owner = %s
                  AND id_form = %s
                  AND id_prerequis = %s
                """,
                (titre, r1, r2, r3, item_ordre, oid, id_form, pid),
            )

            if cur.rowcount > 0:
                ordre += 1
                continue

        cur.execute(
            """
            INSERT INTO public.tbl_fiche_formation_prerequis
              (
                id_prerequis,
                id_owner,
                id_form,
                titre,
                r1,
                r2,
                r3,
                ordre_affichage,
                archive,
                date_creation,
                date_modification
              )
            VALUES
              (
                %s, %s, %s,
                %s, %s, %s, %s,
                %s,
                FALSE,
                NOW(),
                NOW()
              )
            """,
            (
                str(uuid.uuid4()),
                oid,
                id_form,
                titre,
                r1,
                r2,
                r3,
                item_ordre,
            ),
        )

        ordre += 1

def _fetch_form_detail(cur, oid: str, id_form: str) -> dict:
    cur.execute(
        """
        SELECT
          ff.id_form,
          ff.code,
          ff.titre,
          ff.fournisseur_formation,
          fo.nom AS fournisseur_nom,
          ff.type_formation,
          ff.obs_type_form,
          ff.duree,
          ff.objectifs,
          ff.public_cible,
          ff.presentation,
          ff.modalites,
          ff.methode_peda,
          ff.methode_eval,
          ff.competences_stagiaires,
          ff.competences_formateurs,
          ff.attestation_specifique,
          ff.domaine,
          df.titre_court AS domaine_titre_court,
          df.titre AS domaine_titre,
          ff.tarif_mini,
          ff.etat,
          COALESCE(ff.masque, FALSE) AS masque,
          COALESCE(ff.archive, FALSE) AS archive,
          ff.chemin_sharepoint,
          ff.date_creation,
          ff.date_modification
        FROM public.tbl_fiche_formation ff
        LEFT JOIN public.tbl_fournisseur fo
          ON fo.id_owner = ff.id_owner
         AND fo.id_fourn = ff.fournisseur_formation
         AND COALESCE(fo.archive, FALSE) = FALSE
         AND COALESCE(fo.masque, FALSE) = FALSE
        LEFT JOIN public.tbl_domaine_formation df
          ON df.id_domaine_formation = ff.domaine
         AND COALESCE(df.masque, FALSE) = FALSE
        WHERE ff.id_owner = %s
          AND ff.id_form = %s
        LIMIT 1
        """,
        (oid, id_form),
    )

    form = cur.fetchone()
    if not form:
        raise HTTPException(status_code=404, detail="Formation introuvable.")

    modalites_ids = _json_list(form.get("modalites"))
    methode_peda_ids = _json_list(form.get("methode_peda"))
    methode_eval_ids = _json_list(form.get("methode_eval"))
    comp_stag_ids = _json_list(form.get("competences_stagiaires"))
    comp_form_ids = _json_list(form.get("competences_formateurs"))

    form["modalites_ids"] = modalites_ids
    form["methode_peda_ids"] = methode_peda_ids
    form["methode_eval_ids"] = methode_eval_ids
    form["competences_stagiaires_ids"] = comp_stag_ids
    form["competences_formateurs_ids"] = comp_form_ids

    form["modalites_items"] = _resolve_labels(cur, "tbl_mod_form", "id_mod_form", modalites_ids)
    form["methode_peda_items"] = _resolve_labels(cur, "tbl_met_peda", "id_met_peda", methode_peda_ids)
    form["methode_eval_items"] = _resolve_labels(cur, "tbl_met_eval", "id_met_eval", methode_eval_ids)
    form["competences_stagiaires_items"] = _resolve_competences(cur, oid, comp_stag_ids)
    form["competences_formateurs_items"] = _resolve_competences(cur, oid, comp_form_ids)

    cur.execute(
        """
        SELECT
          id_prerequis,
          titre,
          r1,
          r2,
          r3,
          ordre_affichage
        FROM public.tbl_fiche_formation_prerequis
        WHERE id_owner = %s
          AND id_form = %s
          AND COALESCE(archive, FALSE) = FALSE
        ORDER BY COALESCE(ordre_affichage, 999999), titre
        """,
        (oid, id_form),
    )
    form["prerequis"] = cur.fetchall() or []

    cur.execute(
        """
        SELECT
          l.id_ligne_contenu,
          l.titre_sequence,
          l.objectif,
          l.contenu,
          l.id_competence,
          l.competences_liees,
          c.code AS competence_code,
          c.intitule AS competence_intitule,
          l.position
        FROM public.tbl_contenu_ligne l
        LEFT JOIN public.tbl_competence c
          ON c.id_owner = l.id_owner
         AND c.id_comp = l.id_competence
         AND COALESCE(c.masque, FALSE) = FALSE
        WHERE l.id_owner = %s
          AND l.id_form = %s
          AND COALESCE(l.archive, FALSE) = FALSE
        ORDER BY COALESCE(l.position, 999999), l.titre_sequence
        """,
        (oid, id_form),
    )
    contenus = cur.fetchall() or []

    for l in contenus:
        comp_ids = _normalize_competence_ids(l.get("competences_liees"))
        if not comp_ids and l.get("id_competence"):
            comp_ids = [str(l.get("id_competence")).strip()]

        l["competences_liees_ids"] = comp_ids
        l["competences_liees_items"] = _resolve_competences(cur, oid, comp_ids)

    form["contenus"] = contenus

    cur.execute(
        """
        SELECT
          p.id_plan_peda,
          p.codification,
          p.titre,
          p.commentaire,
          p.modalite_generale,
          p.chemin_sharepoint,
          p.date_creation,
          COALESCE(
            SUM(
              CASE
                WHEN trim(COALESCE(b.duree::text, '')) ~ '^[0-9]+([\\.,][0-9]+)?$'
                THEN replace(trim(b.duree::text), ',', '.')::numeric
                ELSE 0
              END
            ),
            0
          ) AS duree_totale,
          COUNT(DISTINCT b.id_bloc_peda) AS nb_blocs
        FROM public.tbl_plan_pedagogique p
        LEFT JOIN public.tbl_bloc_pedagogique b
          ON b.id_owner = p.id_owner
         AND b.id_plan_peda = p.id_plan_peda
         AND COALESCE(b.archive, FALSE) = FALSE
        WHERE p.id_owner = %s
          AND p.id_form = %s
          AND COALESCE(p.archive, FALSE) = FALSE
        GROUP BY
          p.id_plan_peda,
          p.codification,
          p.titre,
          p.commentaire,
          p.modalite_generale,
          p.chemin_sharepoint,
          p.date_creation
        ORDER BY lower(COALESCE(p.codification, '')), lower(COALESCE(p.titre, ''))
        """,
        (oid, id_form),
    )
    plans = cur.fetchall() or []

    for p in plans:
        cur.execute(
            """
            SELECT
              b.id_bloc_peda,
              b.titre,
              b.objectif,
              b.duree,
              b.modalite_intervention,
              b.observations,
              b.position,
              COUNT(s.id_sequence_bloc) AS nb_sequences
            FROM public.tbl_bloc_pedagogique b
            LEFT JOIN public.tbl_sequence_bloc_pedagogique s
              ON s.id_owner = b.id_owner
             AND s.id_bloc_peda = b.id_bloc_peda
             AND COALESCE(s.archive, FALSE) = FALSE
            WHERE b.id_owner = %s
              AND b.id_plan_peda = %s
              AND COALESCE(b.archive, FALSE) = FALSE
            GROUP BY
              b.id_bloc_peda,
              b.titre,
              b.objectif,
              b.duree,
              b.modalite_intervention,
              b.observations,
              b.position
            ORDER BY COALESCE(b.position, 999999), lower(COALESCE(b.titre, ''))
            """,
            (oid, p.get("id_plan_peda")),
        )
        p["blocs"] = cur.fetchall() or []

    form["plans"] = plans
    return form


# ======================================================
# Import document formation
# ======================================================

def _doc_clean_text(value: Any, max_len: int = 22000) -> str:
    txt = str(value or "").replace("\x00", " ")
    txt = re.sub(r"\r\n?", "\n", txt)
    txt = re.sub(r"[ \t]+", " ", txt)
    txt = re.sub(r"\n{3,}", "\n\n", txt).strip()

    if len(txt) > max_len:
        txt = txt[:max_len].rsplit(" ", 1)[0].strip()

    return txt


async def _extract_training_document_text(upload: UploadFile) -> str:
    filename = (upload.filename or "").strip()
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    raw = await upload.read()

    if not raw:
        raise HTTPException(status_code=400, detail="Document vide.")

    if len(raw) > 10 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="Document trop volumineux (10 Mo maximum).")

    if ext == "pdf":
        if PdfReader is None:
            raise HTTPException(status_code=500, detail="Lecture PDF indisponible côté serveur.")

        try:
            reader = PdfReader(BytesIO(raw))
            parts = []

            for page in reader.pages[:40]:
                try:
                    parts.append(page.extract_text() or "")
                except Exception:
                    pass

            txt = _doc_clean_text("\n".join(parts))

            if not txt:
                raise HTTPException(status_code=400, detail="Aucun texte exploitable détecté dans le PDF.")

            return txt

        except HTTPException:
            raise
        except Exception:
            raise HTTPException(status_code=400, detail="Impossible de lire le document PDF.")

    if ext == "docx":
        if DocxDocument is None:
            raise HTTPException(status_code=500, detail="Lecture DOCX indisponible côté serveur.")

        try:
            doc = DocxDocument(BytesIO(raw))
            parts = [p.text for p in doc.paragraphs if p.text]

            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join((cell.text or "").strip() for cell in row.cells))

            txt = _doc_clean_text("\n".join(parts))

            if not txt:
                raise HTTPException(status_code=400, detail="Aucun texte exploitable détecté dans le DOCX.")

            return txt

        except HTTPException:
            raise
        except Exception:
            raise HTTPException(status_code=400, detail="Impossible de lire le document DOCX.")

    if ext == "pptx":
        if PptxPresentation is None:
            raise HTTPException(status_code=500, detail="Lecture PPTX indisponible côté serveur.")

        try:
            prs = PptxPresentation(BytesIO(raw))
            parts = []

            for slide in prs.slides[:80]:
                for shape in slide.shapes:
                    txt = getattr(shape, "text", "") or ""
                    if txt.strip():
                        parts.append(txt)

            txt = _doc_clean_text("\n".join(parts))

            if not txt:
                raise HTTPException(status_code=400, detail="Aucun texte exploitable détecté dans le PPTX.")

            return txt

        except HTTPException:
            raise
        except Exception:
            raise HTTPException(status_code=400, detail="Impossible de lire le document PPTX.")

    raise HTTPException(status_code=400, detail="Format non pris en charge. Formats acceptés : PDF, DOCX ou PPTX.")


def _norm_match_text(value: Any) -> str:
    s = str(value or "").lower()
    s = unicodedata.normalize("NFKD", s)
    s = "".join(ch for ch in s if not unicodedata.combining(ch))
    s = re.sub(r"[^a-z0-9]+", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _token_set(value: Any) -> set:
    stop = {
        "de", "des", "du", "la", "le", "les", "un", "une", "et", "ou", "en",
        "a", "au", "aux", "dans", "sur", "pour", "par", "avec", "sans", "d",
        "l", "formation", "competence", "competences", "capable", "etre", "être"
    }
    return {x for x in _norm_match_text(value).split() if len(x) > 2 and x not in stop}


def _match_score(a: Any, b: Any) -> int:
    aa = _norm_match_text(a)
    bb = _norm_match_text(b)

    if not aa or not bb:
        return 0

    ratio = SequenceMatcher(None, aa, bb).ratio()

    ta = _token_set(aa)
    tb = _token_set(bb)

    if ta and tb:
        overlap = len(ta & tb) / max(1, len(ta | tb))
    else:
        overlap = 0

    score = int(round((ratio * 55) + (overlap * 45)))
    return max(0, min(100, score))


def _fetch_import_catalogue(cur, oid: str) -> list:
    cur.execute(
        """
        SELECT
          c.id_comp,
          c.code,
          c.intitule,
          COALESCE(c.description, '') AS description,
          COALESCE(c.niveaua, '') AS niveaua,
          COALESCE(c.niveaub, '') AS niveaub,
          COALESCE(c.niveauc, '') AS niveauc,
          c.domaine,
          dc.titre_court AS domaine_titre_court,
          dc.titre AS domaine_titre,
          dc.couleur AS domaine_couleur,
          c.etat
        FROM public.tbl_competence c
        LEFT JOIN public.tbl_domaine_competence dc
          ON dc.id_domaine_competence = c.domaine
         AND COALESCE(dc.masque, FALSE) = FALSE
        WHERE c.id_owner = %s
          AND COALESCE(c.masque, FALSE) = FALSE
          AND COALESCE(c.etat, 'active') IN ('active', 'valide', 'à valider')
        ORDER BY lower(COALESCE(c.code, '')), lower(COALESCE(c.intitule, ''))
        """,
        (oid,),
    )

    return cur.fetchall() or []


def _match_import_competences(cur, oid: str, labels: list) -> list:
    catalogue = _fetch_import_catalogue(cur, oid)
    out = []

    for raw in labels or []:
        label = str(raw or "").strip()
        if not label:
            continue

        scored = []

        for c in catalogue:
            blob = " ".join([
                c.get("intitule") or "",
                c.get("description") or "",
                c.get("niveaua") or "",
                c.get("niveaub") or "",
                c.get("niveauc") or "",
                c.get("domaine_titre_court") or "",
                c.get("domaine_titre") or "",
            ])

            score_title = _match_score(label, c.get("intitule") or "")
            score_full = _match_score(label, blob)
            score = max(score_title, int(round((score_title * 0.70) + (score_full * 0.30))))

            if score >= 50:
                scored.append({
                    "id_comp": c.get("id_comp"),
                    "code": c.get("code"),
                    "intitule": c.get("intitule"),
                    "domaine": c.get("domaine"),
                    "domaine_titre_court": c.get("domaine_titre_court"),
                    "domaine_titre": c.get("domaine_titre"),
                    "domaine_couleur": c.get("domaine_couleur"),
                    "score": score,
                })

        scored.sort(key=lambda x: x.get("score") or 0, reverse=True)
        matches = scored[:3]
        best_score = matches[0].get("score", 0) if matches else 0

        if best_score >= 82:
            status = "recommandé"
            selected_id = matches[0].get("id_comp")
        elif best_score >= 66:
            status = "approchant"
            selected_id = None
        elif best_score >= 50:
            status = "à vérifier"
            selected_id = None
        else:
            status = "à créer"
            selected_id = None

        out.append({
            "source": label,
            "selected_id": selected_id,
            "status": status,
            "matches": matches,
        })

    return out


def _find_ref_ids_by_labels(rows: list, id_key: str, labels: list) -> list:
    wanted = [str(x or "").strip() for x in (labels or []) if str(x or "").strip()]
    if not wanted:
        return []

    out = []
    seen = set()

    for label in wanted:
        scored = []

        for r in rows or []:
            rid = str(r.get(id_key) or "").strip()
            if not rid:
                continue

            ref_label = " ".join([
                str(r.get("titre") or ""),
                str(r.get("titre_court") or ""),
                str(r.get("description") or ""),
            ])

            score = _match_score(label, ref_label)
            if score >= 65:
                scored.append((score, rid))

        scored.sort(reverse=True)

        if scored:
            rid = scored[0][1]
            if rid not in seen:
                seen.add(rid)
                out.append(rid)

    return out


def _find_domaine_formation_id(cur, labels: list) -> Optional[str]:
    clean = [str(x or "").strip() for x in (labels or []) if str(x or "").strip()]
    if not clean:
        return None

    cur.execute(
        """
        SELECT
          id_domaine_formation,
          titre,
          titre_court,
          description
        FROM public.tbl_domaine_formation
        WHERE COALESCE(masque, FALSE) = FALSE
        """
    )

    rows = cur.fetchall() or []
    best = None

    for label in clean:
        for r in rows:
            blob = " ".join([
                str(r.get("titre") or ""),
                str(r.get("titre_court") or ""),
                str(r.get("description") or ""),
            ])
            score = _match_score(label, blob)

            if best is None or score > best[0]:
                best = (score, r.get("id_domaine_formation"))

    if best and best[0] >= 65:
        return best[1]

    return None

def _is_generic_prereq_title(value: Any) -> bool:
    s = _norm_match_text(value)

    if not s:
        return True

    generic = {
        "prerequis",
        "pre requis",
        "prerequis formation",
        "conditions prealables",
        "conditions d acces",
        "avant formation",
        "niveau requis",
        "public concerne",
    }

    return s in generic or len(s) < 4


def _looks_like_answer_label(value: Any) -> bool:
    s = _norm_match_text(value)

    if not s:
        return True

    labels = {
        "oui",
        "non",
        "partiellement",
        "je maitrise",
        "je ne maitrise pas",
        "je maitrise partiellement",
        "acquis",
        "non acquis",
        "en cours d acquisition",
        "besoin d assistance",
        "j ai besoin d assistance",
    }

    return s in labels or len(s.split()) <= 3


def _normalize_import_prerequis(raw_items: Any) -> list:
    """
    Normalise les prérequis extraits d'un document.

    Objectif :
    - éviter que l'IA place plusieurs prérequis dans r1/r2/r3 ;
    - conserver un item = un prérequis ;
    - réponses par défaut compatibles Oui / Non / optionnel.
    """
    items = raw_items if isinstance(raw_items, list) else []
    out = []
    seen = set()

    def add_item(titre: Any, r1: Any = "Oui", r2: Any = "Non", r3: Any = ""):
        clean_titre = _clean_text(titre, 800)
        if not clean_titre:
            return

        key = _norm_match_text(clean_titre)
        if not key or key in seen:
            return

        seen.add(key)

        out.append({
            "id_prerequis": None,
            "titre": clean_titre,
            "r1": _clean_text(r1 or "Oui", 300),
            "r2": _clean_text(r2 or "Non", 300),
            "r3": _clean_text(r3 or "", 300),
            "ordre_affichage": len(out) + 1,
        })

    for raw in items:
        if not isinstance(raw, dict):
            add_item(raw)
            continue

        titre = _clean_text(raw.get("titre"), 800)
        r1 = _clean_text(raw.get("r1"), 300)
        r2 = _clean_text(raw.get("r2"), 300)
        r3 = _clean_text(raw.get("r3"), 300)

        responses = [x for x in [r1, r2, r3] if x]

        # Cas typique du bug :
        # titre = "Prérequis", r1/r2/r3 = les vrais prérequis.
        if _is_generic_prereq_title(titre):
            for resp in responses:
                if resp and not _looks_like_answer_label(resp):
                    add_item(resp)
            continue

        # Cas normal : titre = le prérequis, réponses = labels d'auto-positionnement.
        clean_r1 = r1 if _looks_like_answer_label(r1) else "Oui"
        clean_r2 = r2 if _looks_like_answer_label(r2) else "Non"
        clean_r3 = r3 if _looks_like_answer_label(r3) else ""

        add_item(titre, clean_r1, clean_r2, clean_r3)

    return out

def _build_import_ai_schema() -> dict:
    return {
        "type": "object",
        "additionalProperties": False,
        "required": [
            "titre", "presentation", "public_cible", "objectifs", "type_formation",
            "obs_type_form", "duree", "tarif_mini", "domaines_probables",
            "modalites", "methodes_peda", "methodes_eval", "prerequis",
            "competences_stagiaires", "competences_formateurs", "contenus"
        ],
        "properties": {
            "titre": {"type": "string"},
            "presentation": {"type": "string"},
            "public_cible": {"type": "string"},
            "objectifs": {"type": "string"},
            "type_formation": {"type": "string"},
            "obs_type_form": {"type": ["string", "null"]},
            "duree": {"type": ["number", "null"]},
            "tarif_mini": {"type": ["number", "null"]},
            "domaines_probables": {"type": "array", "items": {"type": "string"}},
            "modalites": {"type": "array", "items": {"type": "string"}},
            "methodes_peda": {"type": "array", "items": {"type": "string"}},
            "methodes_eval": {"type": "array", "items": {"type": "string"}},
            "prerequis": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "required": ["titre", "r1", "r2", "r3"],
                    "properties": {
                        "titre": {"type": "string"},
                        "r1": {"type": "string"},
                        "r2": {"type": "string"},
                        "r3": {"type": "string"},
                    },
                },
            },
            "competences_stagiaires": {"type": "array", "items": {"type": "string"}},
            "competences_formateurs": {"type": "array", "items": {"type": "string"}},
            "contenus": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "required": ["titre_sequence", "objectif", "contenu", "competences_sources"],
                    "properties": {
                        "titre_sequence": {"type": "string"},
                        "objectif": {"type": ["string", "null"]},
                        "contenu": {"type": "string"},
                        "competences_sources": {"type": "array", "items": {"type": "string"}},
                    },
                },
            },
        },
    }


def _analyse_import_document_with_ai(doc_text: str, filename: str) -> dict:
    if OpenAI is None:
        raise HTTPException(status_code=500, detail="Lib OpenAI manquante côté serveur.")

    api_key = (os.getenv("OPENAI_API_KEY") or "").strip()
    if not api_key:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY non configurée.")

    model = (os.getenv("OPENAI_MODEL_FORM_IMPORT") or "gpt-4o-mini").strip()

    system_prompt = (
        "Tu analyses un document de formation importé dans Novoskill Learn. "
        "Tu dois extraire uniquement les informations présentes ou fortement déductibles du document. "
        "Tu ne dois pas inventer une formation complète si le document est incomplet. "
        "Tu dois renvoyer un JSON strict conforme au schéma. "
        "Pour les compétences, distingue les compétences visées pour les stagiaires et les compétences requises pour le formateur. "
        "Pour les contenus, découpe en lignes de contenu pédagogiques réutilisables, pas en planning horaire. "
        "Pour les prérequis, règle stricte : chaque élément de la liste prerequis représente UN SEUL prérequis. "
        "Le champ titre contient le prérequis lui-même. "
        "Les champs r1, r2, r3 sont uniquement des réponses d’auto-positionnement, jamais des prérequis. "
        "Pour un prérequis oui/non, utilise r1='Oui', r2='Non', r3=''. "
        "Le type_formation doit être l'une des valeurs: Certifiante, Diplomante, Non Certifiante. "
        "Si une information est absente, renvoie une chaîne vide, null ou une liste vide."
    )

    user_prompt = (
        f"Nom du fichier: {filename}\n\n"
        f"Texte extrait du document:\n{doc_text}"
    )

    client = OpenAI(api_key=api_key)

    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.15,
        max_tokens=3200,
        response_format={
            "type": "json_schema",
            "json_schema": {
                "name": "learn_formation_import",
                "schema": _build_import_ai_schema(),
                "strict": True,
            },
        },
    )

    raw = (resp.choices[0].message.content or "").strip()
    if not raw:
        raise HTTPException(status_code=500, detail="Analyse IA vide.")

    try:
        return json.loads(raw)
    except Exception:
        raise HTTPException(status_code=500, detail="Analyse IA illisible.")

def _build_generation_ai_schema() -> dict:
    return {
        "type": "object",
        "additionalProperties": False,
        "required": [
            "titre", "presentation", "objectif_pedagogique", "public_cible",
            "type_formation", "obs_type_form", "duree_recommandee", "duree_statut",
            "duree_justification", "domaines_probables", "methodes_peda",
            "methodes_eval", "prerequis", "competences_stagiaires",
            "competences_formateurs", "contenus", "rapport_ia"
        ],
        "properties": {
            "titre": {"type": "string"},
            "presentation": {"type": "string"},
            "objectif_pedagogique": {"type": "string"},
            "public_cible": {"type": "string"},
            "type_formation": {"type": "string"},
            "obs_type_form": {"type": ["string", "null"]},
            "duree_recommandee": {"type": ["number", "null"]},
            "duree_statut": {"type": "string"},
            "duree_justification": {"type": "string"},
            "domaines_probables": {"type": "array", "items": {"type": "string"}},
            "methodes_peda": {"type": "array", "items": {"type": "string"}},
            "methodes_eval": {"type": "array", "items": {"type": "string"}},
            "prerequis": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "required": ["titre", "r1", "r2", "r3"],
                    "properties": {
                        "titre": {"type": "string"},
                        "r1": {"type": "string"},
                        "r2": {"type": "string"},
                        "r3": {"type": "string"},
                    },
                },
            },
            "competences_stagiaires": {"type": "array", "items": {"type": "string"}},
            "competences_formateurs": {"type": "array", "items": {"type": "string"}},
            "contenus": {
                "type": "array",
                "items": {
                    "type": "object",
                    "additionalProperties": False,
                    "required": ["titre_sequence", "objectif", "contenu", "competences_sources"],
                    "properties": {
                        "titre_sequence": {"type": "string"},
                        "objectif": {"type": ["string", "null"]},
                        "contenu": {"type": "string"},
                        "competences_sources": {"type": "array", "items": {"type": "string"}},
                    },
                },
            },
            "rapport_ia": {"type": "string"},
        },
    }


def _analyse_generate_formation_with_ai(
    objectif: str,
    contexte: str,
    public_vise: str,
    duree_souhaitee: Optional[float],
    contraintes: str,
    documents_text: str,
) -> dict:
    if OpenAI is None:
        raise HTTPException(status_code=500, detail="Lib OpenAI manquante côté serveur.")

    api_key = (os.getenv("OPENAI_API_KEY") or "").strip()
    if not api_key:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY non configurée.")

    model = (os.getenv("OPENAI_MODEL_FORM_GENERATE") or os.getenv("OPENAI_MODEL_FORM_IMPORT") or "gpt-4o-mini").strip()

    duree_line = "Durée souhaitée non renseignée. Propose une durée réaliste et justifie-la."
    if duree_souhaitee is not None:
        duree_line = f"Durée souhaitée par l'utilisateur : {duree_souhaitee} heure(s). Tu peux la challenger et proposer une durée plus cohérente si nécessaire."

    system_prompt = (
        "Tu génères une fiche formation structurée pour Novoskill Learn. "
        "Règles impératives : zéro marketing, rédaction opérationnelle, pédagogique, sobre et exploitable. "
        "Le titre commence par un verbe d'action. "
        "La présentation est un texte rédigé avec une bonne syntaxe, sans liste. "
        "L'objectif pédagogique est la finalité globale de la formation et doit être formulé dans l'esprit : "
        "À la fin de la formation, le stagiaire/l'apprenant sera capable de... "
        "Ne confonds pas objectif pédagogique et compétences visées. Les compétences visées sont des capacités opérationnelles observables. "
        "Les contenus sont des briques réutilisables indépendantes de la modalité. "
        "Ne génère jamais de plan pédagogique, ni de déroulé jour par jour. "
        "Ne demande pas et ne déduis pas une modalité de réalisation : la modalité sera traitée dans le plan pédagogique. "
        "Les prérequis doivent être évaluables. Chaque prérequis est un élément distinct avec titre + réponses Oui/Non et r3 optionnelle. "
        "Les compétences formateur doivent inclure des compétences d'animation/évaluation/apprentissage et des compétences métier liées au contenu ; "
        "on considère que le formateur possède un niveau avancé ou expert. "
        "Si un contexte est fourni, rends la formation spécifique à ce contexte. Si aucun contexte n'est fourni, génère une formation générique et réutilisable. "
        "Le rapport IA final justifie la compréhension du besoin, les compétences, les contenus, la durée et les points de vigilance. "
        "Le type_formation doit être l'une des valeurs : Certifiante, Diplomante, Non Certifiante. "
        "Renvoie uniquement un JSON strict conforme au schéma."
    )

    user_prompt = (
        f"Objectif demandé :\n{objectif}\n\n"
        f"Contexte optionnel :\n{contexte or 'Non renseigné'}\n\n"
        f"Public visé optionnel :\n{public_vise or 'Non renseigné'}\n\n"
        f"{duree_line}\n\n"
        f"Contraintes éventuelles :\n{contraintes or 'Non renseigné'}\n\n"
        f"Documents de référence extraits :\n{documents_text or 'Aucun document fourni.'}"
    )

    client = OpenAI(api_key=api_key)

    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.2,
        max_tokens=4200,
        response_format={
            "type": "json_schema",
            "json_schema": {
                "name": "learn_formation_generation",
                "schema": _build_generation_ai_schema(),
                "strict": True,
            },
        },
    )

    raw = (resp.choices[0].message.content or "").strip()
    if not raw:
        raise HTTPException(status_code=500, detail="Génération IA vide.")

    try:
        return json.loads(raw)
    except Exception:
        raise HTTPException(status_code=500, detail="Génération IA illisible.")


def _format_generation_report(draft: dict, duree_souhaitee: Optional[float]) -> str:
    report = _clean_text(draft.get("rapport_ia"), 8000)
    if report:
        return report

    duree_txt = "Non renseignée" if duree_souhaitee is None else f"{duree_souhaitee} h"
    rec_txt = draft.get("duree_recommandee") or "Non estimée"

    return _clean_text(
        "\n".join([
            "Rapport de génération IA - Novoskill Learn",
            "",
            f"Formation proposée : {draft.get('titre') or '—'}",
            "",
            f"Durée demandée : {duree_txt}",
            f"Durée recommandée : {rec_txt} h",
            f"Analyse durée : {draft.get('duree_justification') or '—'}",
            "",
            "Points de vigilance : vérifier les compétences proposées à création avant enregistrement définitif.",
        ]),
        8000,
    )

# ======================================================
# Models
# ======================================================

class FormationPrerequisPayload(BaseModel):
    id_prerequis: Optional[str] = None
    titre: Optional[str] = None
    r1: Optional[str] = None
    r2: Optional[str] = None
    r3: Optional[str] = None
    ordre_affichage: Optional[int] = None


class FormationPayload(BaseModel):
    titre: str
    fournisseur_formation: Optional[str] = None
    type_formation: Optional[str] = None
    obs_type_form: Optional[str] = None
    duree: Optional[Any] = None
    objectifs: Optional[str] = None
    public_cible: Optional[str] = None
    presentation: Optional[str] = None
    modalites: Optional[list] = None
    methode_peda: Optional[list] = None
    methode_eval: Optional[list] = None
    competences_stagiaires: Optional[list] = None
    competences_formateurs: Optional[list] = None
    prerequis: Optional[list[FormationPrerequisPayload]] = None
    attestation_specifique: Optional[str] = None
    domaine: Optional[str] = None
    tarif_mini: Optional[Any] = None
    etat: Optional[str] = None


class FormationUpdatePayload(BaseModel):
    titre: Optional[str] = None
    fournisseur_formation: Optional[str] = None
    type_formation: Optional[str] = None
    obs_type_form: Optional[str] = None
    duree: Optional[Any] = None
    objectifs: Optional[str] = None
    public_cible: Optional[str] = None
    presentation: Optional[str] = None
    modalites: Optional[list] = None
    methode_peda: Optional[list] = None
    methode_eval: Optional[list] = None
    competences_stagiaires: Optional[list] = None
    competences_formateurs: Optional[list] = None
    prerequis: Optional[list[FormationPrerequisPayload]] = None
    attestation_specifique: Optional[str] = None
    domaine: Optional[str] = None
    tarif_mini: Optional[Any] = None
    etat: Optional[str] = None

class FormationContenuPayload(BaseModel):
    titre_sequence: str
    objectif: Optional[str] = None
    contenu: Optional[str] = None
    competences_liees: Optional[list] = None
    position: Optional[int] = None


class FormationContenuReorderPayload(BaseModel):
    items: list[str]

class FormationPlanBlocPayload(BaseModel):
    titre: Optional[str] = None
    duree: Optional[Any] = None
    modalite_intervention: Optional[str] = None
    objectif: Optional[str] = None
    observations: Optional[str] = None
    contenus: Optional[list] = None
    position: Optional[int] = None


class FormationPlanPayload(BaseModel):
    titre: str
    modalite_generale: Optional[str] = None
    commentaire: Optional[str] = None
    blocs: Optional[list[FormationPlanBlocPayload]] = None

# ======================================================
# Référentiels
# ======================================================

@router.get("/learn/formations/{id_effectif}/context")
def learn_formations_context(id_effectif: str, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    with get_conn() as conn:
        with conn.cursor(row_factory=dict_row) as cur:
            profile = _learn_require_profile(cur, u, id_effectif)

    return {
        "id_owner": profile.get("id_owner"),
        "id_effectif": profile.get("id_effectif"),
        "role_code": profile.get("role_code"),
        "role_label": profile.get("role_label"),
        "can_edit": _role_rank(profile.get("role_code")) >= 2,
    }


@router.post("/learn/formations/{id_effectif}/import_document")
async def learn_formations_import_document(
    id_effectif: str,
    request: Request,
    document: UploadFile = File(...),
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        filename = (document.filename or "").strip()
        if not filename:
            raise HTTPException(status_code=400, detail="Nom de fichier manquant.")

        doc_text = await _extract_training_document_text(document)

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                draft = _analyse_import_document_with_ai(doc_text, filename)

                cur.execute(
                    """
                    SELECT id_mod_form, titre, titre_court, description
                    FROM public.tbl_mod_form
                    WHERE COALESCE(masque, FALSE) = FALSE
                    """
                )
                modalites_rows = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT id_met_peda, titre, titre_court, description
                    FROM public.tbl_met_peda
                    WHERE COALESCE(masque, FALSE) = FALSE
                    """
                )
                peda_rows = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT id_met_eval, titre, titre_court, description
                    FROM public.tbl_met_eval
                    WHERE COALESCE(masque, FALSE) = FALSE
                    """
                )
                eval_rows = cur.fetchall() or []

                domaine_id = _find_domaine_formation_id(cur, draft.get("domaines_probables") or [])

                comp_stag = _match_import_competences(cur, oid, draft.get("competences_stagiaires") or [])
                comp_form = _match_import_competences(cur, oid, draft.get("competences_formateurs") or [])

        prerequis = _normalize_import_prerequis(draft.get("prerequis") or [])

        contenus = []
        for c in draft.get("contenus") or []:
            titre = _clean_text(c.get("titre_sequence"), 500)
            detail = _clean_text(c.get("contenu"), 8000)

            if not titre and not detail:
                continue

            contenus.append({
                "titre_sequence": titre or "Contenu",
                "objectif": _clean_text(c.get("objectif"), 1200),
                "contenu": detail,
                "competences_sources": [
                    str(x or "").strip()
                    for x in (c.get("competences_sources") or [])
                    if str(x or "").strip()
                ],
                "competences_liees": [],
            })

        out = {
            "filename": filename,
            "titre": _clean_text(draft.get("titre"), 500),
            "presentation": _clean_text(draft.get("presentation"), 6000),
            "public_cible": _clean_text(draft.get("public_cible"), 3000),
            "objectifs": _clean_text(draft.get("objectifs"), 5000),
            "type_formation": _normalize_type_formation(draft.get("type_formation")),
            "obs_type_form": _clean_text(draft.get("obs_type_form"), 500),
            "duree": _safe_float(draft.get("duree")),
            "tarif_mini": _safe_float(draft.get("tarif_mini")),
            "domaine": domaine_id,
            "modalites_ids": _find_ref_ids_by_labels(modalites_rows, "id_mod_form", draft.get("modalites") or []),
            "methode_peda_ids": _find_ref_ids_by_labels(peda_rows, "id_met_peda", draft.get("methodes_peda") or []),
            "methode_eval_ids": _find_ref_ids_by_labels(eval_rows, "id_met_eval", draft.get("methodes_eval") or []),
            "prerequis": prerequis,
            "competences_stagiaires_import": comp_stag,
            "competences_formateurs_import": comp_form,
            "contenus": contenus,
            "raw_text_preview": doc_text[:1800],
        }

        return out

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations import_document error: {e}")

@router.post("/learn/formations/{id_effectif}/generate_ai")
async def learn_formations_generate_ai(
    id_effectif: str,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        form = await request.form()

        objectif = str(form.get("objectif") or "")
        contexte = str(form.get("contexte") or "")
        public_vise = str(form.get("public_vise") or "")
        duree_souhaitee = str(form.get("duree_souhaitee") or "")
        contraintes = str(form.get("contraintes") or "")

        documents = []
        try:
            raw_docs = form.getlist("documents")
        except Exception:
            raw_docs = []

        for item in raw_docs:
            if hasattr(item, "filename") and hasattr(item, "read"):
                if (item.filename or "").strip():
                    documents.append(item)

        obj = _clean_text(objectif, 3000)
        if not obj:
            raise HTTPException(status_code=400, detail="Objectif de formation obligatoire.")

        duree = _safe_float(duree_souhaitee)
        docs_text_parts = []

        for upload in documents or []:
            filename = (upload.filename or "").strip()
            if not filename:
                continue

            txt = await _extract_training_document_text(upload)
            if txt:
                docs_text_parts.append(f"--- Document : {filename} ---\n{txt}")

        docs_text = _doc_clean_text("\n\n".join(docs_text_parts), 36000)

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                draft = _analyse_generate_formation_with_ai(
                    objectif=obj,
                    contexte=_clean_text(contexte, 5000),
                    public_vise=_clean_text(public_vise, 2000),
                    duree_souhaitee=duree,
                    contraintes=_clean_text(contraintes, 4000),
                    documents_text=docs_text,
                )

                cur.execute(
                    """
                    SELECT id_met_peda, titre, titre_court, description
                    FROM public.tbl_met_peda
                    WHERE COALESCE(masque, FALSE) = FALSE
                    """
                )
                peda_rows = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT id_met_eval, titre, titre_court, description
                    FROM public.tbl_met_eval
                    WHERE COALESCE(masque, FALSE) = FALSE
                    """
                )
                eval_rows = cur.fetchall() or []

                domaine_id = _find_domaine_formation_id(cur, draft.get("domaines_probables") or [])
                comp_stag = _match_import_competences(cur, oid, draft.get("competences_stagiaires") or [])
                comp_form = _match_import_competences(cur, oid, draft.get("competences_formateurs") or [])

        prerequis = _normalize_import_prerequis(draft.get("prerequis") or [])

        contenus = []
        for c in draft.get("contenus") or []:
            titre = _clean_text(c.get("titre_sequence"), 500)
            detail = _clean_text(c.get("contenu"), 8000)

            if not titre and not detail:
                continue

            contenus.append({
                "titre_sequence": titre or "Contenu",
                "objectif": _clean_text(c.get("objectif"), 1200),
                "contenu": detail,
                "competences_sources": [
                    str(x or "").strip()
                    for x in (c.get("competences_sources") or [])
                    if str(x or "").strip()
                ],
                "competences_liees": [],
            })

        duree_recommandee = _safe_float(draft.get("duree_recommandee"))

        out = {
            "titre": _clean_text(draft.get("titre"), 500),
            "presentation": _clean_text(draft.get("presentation"), 6000),
            "public_cible": _clean_text(draft.get("public_cible"), 3000),
            "objectifs": _clean_text(draft.get("objectif_pedagogique"), 5000),
            "type_formation": _normalize_type_formation(draft.get("type_formation")),
            "obs_type_form": _clean_text(draft.get("obs_type_form"), 500),
            "duree": duree_recommandee,
            "duree_demandee": duree,
            "duree_recommandee": duree_recommandee,
            "duree_statut": _clean_text(draft.get("duree_statut"), 300),
            "duree_justification": _clean_text(draft.get("duree_justification"), 1200),
            "tarif_mini": None,
            "domaine": domaine_id,
            "modalites_ids": [],
            "methode_peda_ids": _find_ref_ids_by_labels(peda_rows, "id_met_peda", draft.get("methodes_peda") or []),
            "methode_eval_ids": _find_ref_ids_by_labels(eval_rows, "id_met_eval", draft.get("methodes_eval") or []),
            "prerequis": prerequis,
            "competences_stagiaires_import": comp_stag,
            "competences_formateurs_import": comp_form,
            "contenus": contenus,
            "rapport_ia": _format_generation_report(draft, duree),
        }

        return out

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations generate_ai error: {e}")

@router.get("/learn/formations/{id_effectif}/referentiels")
def learn_formations_referentiels(id_effectif: str, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                oid = (profile.get("id_owner") or "").strip()

                cur.execute(
                    """
                    SELECT
                      id_domaine_formation,
                      titre,
                      titre_court,
                      description,
                      ordre_affichage
                    FROM public.tbl_domaine_formation
                    WHERE COALESCE(masque, FALSE) = FALSE
                    ORDER BY COALESCE(ordre_affichage, 999999), lower(COALESCE(titre, titre_court, id_domaine_formation))
                    """
                )
                domaines = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT
                      id_fourn,
                      code,
                      nom
                    FROM public.tbl_fournisseur
                    WHERE id_owner = %s
                      AND COALESCE(archive, FALSE) = FALSE
                      AND COALESCE(masque, FALSE) = FALSE
                      AND COALESCE(formation, FALSE) = TRUE
                    ORDER BY lower(COALESCE(nom, code, id_fourn))
                    """,
                    (oid,),
                )
                fournisseurs = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT
                      id_mod_form,
                      titre,
                      titre_court,
                      description,
                      ordre_affichage
                    FROM public.tbl_mod_form
                    WHERE COALESCE(masque, FALSE) = FALSE
                    ORDER BY COALESCE(ordre_affichage, 999999), lower(COALESCE(titre, titre_court, id_mod_form))
                    """
                )
                modalites = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT
                      id_met_peda,
                      titre,
                      titre_court,
                      description,
                      ordre_affichage
                    FROM public.tbl_met_peda
                    WHERE COALESCE(masque, FALSE) = FALSE
                    ORDER BY COALESCE(ordre_affichage, 999999), lower(COALESCE(titre, titre_court, id_met_peda))
                    """
                )
                methodes_peda = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT
                      id_met_eval,
                      titre,
                      titre_court,
                      description,
                      ordre_affichage
                    FROM public.tbl_met_eval
                    WHERE COALESCE(masque, FALSE) = FALSE
                    ORDER BY COALESCE(ordre_affichage, 999999), lower(COALESCE(titre, titre_court, id_met_eval))
                    """
                )
                methodes_eval = cur.fetchall() or []

                cur.execute(
                    """
                    SELECT
                      c.id_comp,
                      c.code,
                      c.intitule,
                      c.domaine,
                      dc.titre_court AS domaine_titre_court,
                      dc.titre AS domaine_titre,
                      dc.couleur AS domaine_couleur,
                      c.etat
                    FROM public.tbl_competence c
                    LEFT JOIN public.tbl_domaine_competence dc
                      ON dc.id_domaine_competence = c.domaine
                     AND COALESCE(dc.masque, FALSE) = FALSE
                    WHERE c.id_owner = %s
                      AND COALESCE(c.masque, FALSE) = FALSE
                      AND COALESCE(c.etat, 'active') IN ('active', 'valide', 'à valider')
                    ORDER BY lower(COALESCE(c.code, '')), lower(COALESCE(c.intitule, ''))
                    """,
                    (oid,),
                )
                competences = cur.fetchall() or []

        return {
            "domaines": domaines,
            "fournisseurs": fournisseurs,
            "modalites": modalites,
            "methodes_peda": methodes_peda,
            "methodes_eval": methodes_eval,
            "competences": competences,
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations/referentiels error: {e}")


# ======================================================
# Catalogue formations
# ======================================================

@router.get("/learn/formations/{id_effectif}")
def learn_formations_list(
    id_effectif: str,
    request: Request,
    q: str = "",
    show: str = "active",
    domaine: str = "",
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        qq = (q or "").strip()
        sh = (show or "active").strip().lower()
        dom = (domaine or "").strip()

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                oid = (profile.get("id_owner") or "").strip()

                where = ["ff.id_owner = %s"]
                params = [oid]

                if sh == "active":
                    where.append("COALESCE(ff.archive, FALSE) = FALSE")
                    where.append("COALESCE(ff.masque, FALSE) = FALSE")
                    where.append("COALESCE(ff.etat, 'active') = 'active'")
                elif sh == "validation":
                    where.append("COALESCE(ff.archive, FALSE) = FALSE")
                    where.append("COALESCE(ff.masque, FALSE) = FALSE")
                    where.append("COALESCE(ff.etat, '') = 'à valider'")
                elif sh == "archived":
                    where.append("(COALESCE(ff.archive, FALSE) = TRUE OR COALESCE(ff.masque, FALSE) = TRUE)")
                else:
                    where.append("COALESCE(ff.archive, FALSE) = FALSE")
                    where.append("COALESCE(ff.masque, FALSE) = FALSE")

                if dom:
                    where.append("ff.domaine = %s")
                    params.append(dom)

                if qq:
                    like = f"%{qq}%"
                    where.append(
                        """
                        (
                          ff.code ILIKE %s
                          OR ff.titre ILIKE %s
                          OR COALESCE(ff.presentation, '') ILIKE %s
                          OR COALESCE(ff.objectifs, '') ILIKE %s
                          OR COALESCE(df.titre, '') ILIKE %s
                          OR COALESCE(df.titre_court, '') ILIKE %s
                          OR COALESCE(fo.nom, '') ILIKE %s
                        )
                        """
                    )
                    params.extend([like, like, like, like, like, like, like])

                cur.execute(
                    f"""
                    SELECT
                      ff.id_form,
                      ff.code,
                      ff.titre,
                      ff.type_formation,
                      ff.duree,
                      ff.domaine,
                      df.titre_court AS domaine_titre_court,
                      df.titre AS domaine_titre,
                      ff.fournisseur_formation,
                      fo.nom AS fournisseur_nom,
                      ff.tarif_mini,
                      ff.etat,
                      COALESCE(ff.masque, FALSE) AS masque,
                      COALESCE(ff.archive, FALSE) AS archive,
                      COUNT(DISTINCT p.id_plan_peda) AS nb_plans
                    FROM public.tbl_fiche_formation ff
                    LEFT JOIN public.tbl_domaine_formation df
                      ON df.id_domaine_formation = ff.domaine
                     AND COALESCE(df.masque, FALSE) = FALSE
                    LEFT JOIN public.tbl_fournisseur fo
                      ON fo.id_owner = ff.id_owner
                     AND fo.id_fourn = ff.fournisseur_formation
                     AND COALESCE(fo.archive, FALSE) = FALSE
                     AND COALESCE(fo.masque, FALSE) = FALSE
                    LEFT JOIN public.tbl_plan_pedagogique p
                      ON p.id_owner = ff.id_owner
                     AND p.id_form = ff.id_form
                     AND COALESCE(p.archive, FALSE) = FALSE
                    WHERE {" AND ".join(where)}
                    GROUP BY
                      ff.id_form,
                      ff.code,
                      ff.titre,
                      ff.type_formation,
                      ff.duree,
                      ff.domaine,
                      df.titre_court,
                      df.titre,
                      ff.fournisseur_formation,
                      fo.nom,
                      ff.tarif_mini,
                      ff.etat,
                      ff.masque,
                      ff.archive
                    ORDER BY
                      lower(COALESCE(ff.code, '')),
                      lower(COALESCE(ff.titre, ''))
                    """,
                    tuple(params),
                )

                rows = cur.fetchall() or []

        return {"items": rows}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations list error: {e}")


@router.get("/learn/formations/{id_effectif}/{id_form}")
def learn_formation_detail(id_effectif: str, id_form: str, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                oid = (profile.get("id_owner") or "").strip()
                return _fetch_form_detail(cur, oid, fid)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations detail error: {e}")


@router.post("/learn/formations/{id_effectif}")
def learn_formation_create(id_effectif: str, payload: FormationPayload, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        titre = (payload.titre or "").strip()
        if not titre:
            raise HTTPException(status_code=400, detail="Titre obligatoire.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                fid = str(uuid.uuid4())
                code = _next_form_code(cur, oid)

                cur.execute(
                    """
                    INSERT INTO public.tbl_fiche_formation
                      (
                        id_form,
                        id_owner,
                        code,
                        titre,
                        fournisseur_formation,
                        type_formation,
                        obs_type_form,
                        duree,
                        objectifs,
                        public_cible,
                        presentation,
                        modalites,
                        methode_peda,
                        methode_eval,
                        competences_stagiaires,
                        competences_formateurs,
                        attestation_specifique,
                        domaine,
                        tarif_mini,
                        etat,
                        masque,
                        archive,
                        date_creation,
                        date_modification
                      )
                    VALUES
                      (
                        %s, %s, %s, %s,
                        %s, %s, %s, %s,
                        %s, %s, %s,
                        %s::jsonb, %s::jsonb, %s::jsonb,
                        %s::jsonb, %s::jsonb,
                        %s, %s, %s, %s,
                        FALSE, FALSE, NOW(), NOW()
                      )
                    """,
                    (
                        fid,
                        oid,
                        code,
                        titre,
                        (payload.fournisseur_formation or None),
                        _normalize_type_formation(payload.type_formation),
                        _clean_text(payload.obs_type_form),
                        _safe_int(payload.duree),
                        _clean_text(payload.objectifs),
                        _clean_text(payload.public_cible),
                        _clean_text(payload.presentation),
                        _jsonb_param(payload.modalites),
                        _jsonb_param(payload.methode_peda),
                        _jsonb_param(payload.methode_eval),
                        _jsonb_param(payload.competences_stagiaires),
                        _jsonb_param(payload.competences_formateurs),
                        _clean_text(payload.attestation_specifique),
                        (payload.domaine or None),
                        _safe_float(payload.tarif_mini),
                        (payload.etat or "à valider"),
                    ),
                )

                _sync_formation_prerequis(cur, oid, fid, payload.prerequis)

                conn.commit()

        return {"id_form": fid, "code": code}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations create error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}")
def learn_formation_update(id_effectif: str, id_form: str, payload: FormationUpdatePayload, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")

        patch_fields = payload.__fields_set__ or set()
        if not patch_fields:
            return {"ok": True}

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                if not _formation_exists_owner(cur, oid, fid):
                    raise HTTPException(status_code=404, detail="Formation introuvable.")

                cols = []
                vals = []

                if "titre" in patch_fields:
                    titre = (payload.titre or "").strip()
                    if not titre:
                        raise HTTPException(status_code=400, detail="Titre obligatoire.")
                    cols.append("titre = %s")
                    vals.append(titre)

                if "fournisseur_formation" in patch_fields:
                    cols.append("fournisseur_formation = %s")
                    vals.append(payload.fournisseur_formation or None)

                if "type_formation" in patch_fields:
                    cols.append("type_formation = %s")
                    vals.append(_normalize_type_formation(payload.type_formation))

                if "obs_type_form" in patch_fields:
                    cols.append("obs_type_form = %s")
                    vals.append(_clean_text(payload.obs_type_form))

                if "duree" in patch_fields:
                    cols.append("duree = %s")
                    vals.append(_safe_int(payload.duree))

                if "objectifs" in patch_fields:
                    cols.append("objectifs = %s")
                    vals.append(_clean_text(payload.objectifs))

                if "public_cible" in patch_fields:
                    cols.append("public_cible = %s")
                    vals.append(_clean_text(payload.public_cible))

                if "presentation" in patch_fields:
                    cols.append("presentation = %s")
                    vals.append(_clean_text(payload.presentation))

                if "modalites" in patch_fields:
                    cols.append("modalites = %s::jsonb")
                    vals.append(_jsonb_param(payload.modalites))

                if "methode_peda" in patch_fields:
                    cols.append("methode_peda = %s::jsonb")
                    vals.append(_jsonb_param(payload.methode_peda))

                if "methode_eval" in patch_fields:
                    cols.append("methode_eval = %s::jsonb")
                    vals.append(_jsonb_param(payload.methode_eval))

                if "competences_stagiaires" in patch_fields:
                    cols.append("competences_stagiaires = %s::jsonb")
                    vals.append(_jsonb_param(payload.competences_stagiaires))

                if "competences_formateurs" in patch_fields:
                    cols.append("competences_formateurs = %s::jsonb")
                    vals.append(_jsonb_param(payload.competences_formateurs))

                if "attestation_specifique" in patch_fields:
                    cols.append("attestation_specifique = %s")
                    vals.append(_clean_text(payload.attestation_specifique))

                if "domaine" in patch_fields:
                    cols.append("domaine = %s")
                    vals.append(payload.domaine or None)

                if "tarif_mini" in patch_fields:
                    cols.append("tarif_mini = %s")
                    vals.append(_safe_float(payload.tarif_mini))

                if "etat" in patch_fields:
                    cols.append("etat = %s")
                    vals.append((payload.etat or "à valider").strip())

                if "prerequis" in patch_fields:
                    _sync_formation_prerequis(cur, oid, fid, payload.prerequis)

                if cols:
                    cols.append("date_modification = NOW()")
                    vals.extend([fid, oid])

                    cur.execute(
                        f"""
                        UPDATE public.tbl_fiche_formation
                        SET {", ".join(cols)}
                        WHERE id_form = %s
                          AND id_owner = %s
                        """,
                        tuple(vals),
                    )

                    conn.commit()

        return {"ok": True}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations update error: {e}")

@router.post("/learn/formations/{id_effectif}/{id_form}/contenus")
def learn_formation_contenu_create(
    id_effectif: str,
    id_form: str,
    payload: FormationContenuPayload,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        titre = (payload.titre_sequence or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not titre:
            raise HTTPException(status_code=400, detail="Titre du contenu obligatoire.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                if not _formation_exists_owner(cur, oid, fid):
                    raise HTTPException(status_code=404, detail="Formation introuvable.")

                comp_ids = _ensure_competences_owner(cur, oid, payload.competences_liees or [])
                main_comp = comp_ids[0] if comp_ids else None

                if payload.position is not None:
                    position = int(payload.position)
                else:
                    cur.execute(
                        """
                        SELECT COALESCE(MAX(position), 0) + 1 AS next_position
                        FROM public.tbl_contenu_ligne
                        WHERE id_owner = %s
                          AND id_form = %s
                          AND COALESCE(archive, FALSE) = FALSE
                        """,
                        (oid, fid),
                    )
                    position = int((cur.fetchone() or {}).get("next_position") or 1)

                lid = str(uuid.uuid4())

                cur.execute(
                    """
                    INSERT INTO public.tbl_contenu_ligne
                      (
                        id_ligne_contenu,
                        id_owner,
                        id_form,
                        titre_sequence,
                        objectif,
                        contenu,
                        id_competence,
                        competences_liees,
                        position,
                        archive,
                        date_creation,
                        date_modification
                      )
                    VALUES
                      (
                        %s, %s, %s,
                        %s, %s, %s,
                        %s, %s::jsonb,
                        %s,
                        FALSE,
                        NOW(),
                        NOW()
                      )
                    """,
                    (
                        lid,
                        oid,
                        fid,
                        titre,
                        _clean_text(payload.objectif),
                        _clean_text(payload.contenu),
                        main_comp,
                        json.dumps(comp_ids, ensure_ascii=False),
                        position,
                    ),
                )

                _renumber_contenus(cur, oid, fid)
                item = _fetch_contenu_row(cur, oid, fid, lid)

                conn.commit()

        return {"ok": True, "item": item}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations contenu create error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}/contenus/{id_ligne_contenu}")
def learn_formation_contenu_update(
    id_effectif: str,
    id_form: str,
    id_ligne_contenu: str,
    payload: FormationContenuPayload,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        lid = (id_ligne_contenu or "").strip()
        titre = (payload.titre_sequence or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not lid:
            raise HTTPException(status_code=400, detail="id_ligne_contenu manquant.")
        if not titre:
            raise HTTPException(status_code=400, detail="Titre du contenu obligatoire.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                if not _formation_exists_owner(cur, oid, fid):
                    raise HTTPException(status_code=404, detail="Formation introuvable.")

                comp_ids = _ensure_competences_owner(cur, oid, payload.competences_liees or [])
                main_comp = comp_ids[0] if comp_ids else None

                cur.execute(
                    """
                    UPDATE public.tbl_contenu_ligne
                    SET titre_sequence = %s,
                        objectif = %s,
                        contenu = %s,
                        id_competence = %s,
                        competences_liees = %s::jsonb,
                        date_modification = NOW()
                    WHERE id_owner = %s
                      AND id_form = %s
                      AND id_ligne_contenu = %s
                      AND COALESCE(archive, FALSE) = FALSE
                    """,
                    (
                        titre,
                        _clean_text(payload.objectif),
                        _clean_text(payload.contenu),
                        main_comp,
                        json.dumps(comp_ids, ensure_ascii=False),
                        oid,
                        fid,
                        lid,
                    ),
                )

                if cur.rowcount <= 0:
                    raise HTTPException(status_code=404, detail="Contenu introuvable.")

                item = _fetch_contenu_row(cur, oid, fid, lid)

                conn.commit()

        return {"ok": True, "item": item}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations contenu update error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}/contenus/{id_ligne_contenu}/archive")
def learn_formation_contenu_archive(
    id_effectif: str,
    id_form: str,
    id_ligne_contenu: str,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        lid = (id_ligne_contenu or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not lid:
            raise HTTPException(status_code=400, detail="id_ligne_contenu manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                cur.execute(
                    """
                    UPDATE public.tbl_contenu_ligne
                    SET archive = TRUE,
                        date_modification = NOW()
                    WHERE id_owner = %s
                      AND id_form = %s
                      AND id_ligne_contenu = %s
                      AND COALESCE(archive, FALSE) = FALSE
                    """,
                    (oid, fid, lid),
                )

                if cur.rowcount <= 0:
                    raise HTTPException(status_code=404, detail="Contenu introuvable.")

                _renumber_contenus(cur, oid, fid)
                conn.commit()

        return {"ok": True}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations contenu archive error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}/contenus/reorder")
def learn_formation_contenu_reorder(
    id_effectif: str,
    id_form: str,
    payload: FormationContenuReorderPayload,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        ids = [str(x or "").strip() for x in (payload.items or []) if str(x or "").strip()]

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not ids:
            return {"ok": True}

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                if not _formation_exists_owner(cur, oid, fid):
                    raise HTTPException(status_code=404, detail="Formation introuvable.")

                for idx, lid in enumerate(ids, start=1):
                    cur.execute(
                        """
                        UPDATE public.tbl_contenu_ligne
                        SET position = %s,
                            date_modification = NOW()
                        WHERE id_owner = %s
                          AND id_form = %s
                          AND id_ligne_contenu = %s
                          AND COALESCE(archive, FALSE) = FALSE
                        """,
                        (idx, oid, fid, lid),
                    )

                conn.commit()

        return {"ok": True}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations contenu reorder error: {e}")

@router.post("/learn/formations/{id_effectif}/{id_form}/archive")
def learn_formation_archive(id_effectif: str, id_form: str, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                cur.execute(
                    """
                    UPDATE public.tbl_fiche_formation
                    SET archive = TRUE,
                        masque = TRUE,
                        date_modification = NOW()
                    WHERE id_owner = %s
                      AND id_form = %s
                      AND COALESCE(archive, FALSE) = FALSE
                    """,
                    (oid, fid),
                )

                conn.commit()

        return {"ok": True}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations archive error: {e}")


# ======================================================
# PDF socle fiche formation
# ======================================================

def _p(txt: Any, style):
    return Paragraph(html.escape(str(txt or "").replace("\n", "<br/>")), style)


def _bullet_list(items: list, style) -> list:
    out = []
    for it in items:
        if isinstance(it, dict):
            label = it.get("titre") or it.get("intitule") or it.get("nom") or it.get("code") or ""
        else:
            label = str(it or "")
        label = str(label or "").strip()
        if label:
            out.append(_p(f"• {label}", style))
    if not out:
        out.append(_p("—", style))
    return out


def _build_formation_pdf_story(form: dict) -> list:
    styles = build_pdf_styles()

    title_style = ParagraphStyle(
        "LearnFormTitle",
        parent=styles["title"],
        fontName="Helvetica-Bold",
        fontSize=16,
        leading=19,
        textColor=colors.HexColor("#1f2937"),
        spaceAfter=4,
    )

    section_style = ParagraphStyle(
        "LearnFormSection",
        parent=styles["section"],
        fontName="Helvetica-Bold",
        fontSize=11.2,
        leading=13,
        textColor=colors.HexColor("#c2410c"),
        spaceAfter=5,
        spaceBefore=8,
    )

    body_style = ParagraphStyle(
        "LearnFormBody",
        parent=styles["body"],
        fontName="Helvetica",
        fontSize=8.8,
        leading=11.2,
        textColor=colors.HexColor("#1f2937"),
        spaceAfter=2,
    )

    small_style = ParagraphStyle(
        "LearnFormSmall",
        parent=styles["small"],
        fontName="Helvetica",
        fontSize=7.8,
        leading=9.4,
        textColor=colors.HexColor("#6b7280"),
    )

    story = []

    code = form.get("code") or "—"
    titre = form.get("titre") or "Formation"
    domaine = form.get("domaine_titre_court") or form.get("domaine_titre") or "—"
    fournisseur = form.get("fournisseur_nom") or "—"

    story.append(_p("Fiche descriptive de formation", title_style))
    story.append(_p(f"{code} • {titre}", styles["subtitle"]))
    story.append(Spacer(1, 4 * mm))

    meta = [
        [_p("Domaine", small_style), _p(domaine, body_style), _p("Durée", small_style), _p(f"{form.get('duree') or '—'} h", body_style)],
        [_p("Fournisseur", small_style), _p(fournisseur, body_style), _p("Tarif mini", small_style), _p(f"{form.get('tarif_mini') or '—'} € HT", body_style)],
        [_p("Type", small_style), _p(form.get("type_formation") or "—", body_style), _p("État", small_style), _p(form.get("etat") or "—", body_style)],
        [_p("Précision type", small_style), _p(form.get("obs_type_form") or "—", body_style), _p("", small_style), _p("", body_style)],
    ]

    tbl = Table(meta, colWidths=[28 * mm, 62 * mm, 28 * mm, 62 * mm])
    tbl.setStyle(TableStyle([
        ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor("#e5e7eb")),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e5e7eb")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(tbl)

    story.append(_p("Présentation", section_style))
    story.append(_p(form.get("presentation") or "—", body_style))

    story.append(_p("Public cible", section_style))
    story.extend(_bullet_list([x for x in str(form.get("public_cible") or "").split("\n") if x.strip()], body_style))

    story.append(_p("Objectifs de la formation", section_style))
    story.append(_p(form.get("objectifs") or "—", body_style))

    story.append(PageBreak())

    story.append(_p("Modalités de formation", title_style))
    story.append(Spacer(1, 4 * mm))

    story.append(_p("Prérequis d’entrée en formation", section_style))
    prereq = form.get("prerequis") or []
    story.extend(_bullet_list([p.get("titre") for p in prereq], body_style))

    story.append(_p("Objectifs pédagogiques / compétences visées", section_style))
    story.extend(_bullet_list(form.get("competences_stagiaires_items") or [], body_style))

    story.append(_p("Modalités possibles", section_style))
    story.extend(_bullet_list(form.get("modalites_items") or [], body_style))

    story.append(_p("Méthodes pédagogiques", section_style))
    story.extend(_bullet_list(form.get("methode_peda_items") or [], body_style))

    story.append(_p("Méthodes d’évaluation des acquis", section_style))
    story.extend(_bullet_list(form.get("methode_eval_items") or [], body_style))

    story.append(_p("Compétences techniques du formateur", section_style))
    comp_form = form.get("competences_formateurs_items") or []
    rows = [[_p("Code", small_style), _p("Désignation", small_style)]]
    for c in comp_form:
        rows.append([_p(c.get("code") or "—", body_style), _p(c.get("intitule") or "—", body_style)])

    table = Table(rows, colWidths=[30 * mm, 145 * mm])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#fff7ed")),
        ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#e5e7eb")),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e5e7eb")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(table)

    story.append(PageBreak())

    story.append(_p("Contenu de la formation", title_style))
    story.append(Spacer(1, 4 * mm))

    contenus = form.get("contenus") or []
    if not contenus:
        story.append(_p("Aucun contenu détaillé n’est encore rattaché à cette formation.", body_style))
    else:
        for l in contenus:
            story.append(_p(l.get("titre_sequence") or "Séquence", section_style))
            if l.get("objectif"):
                story.append(_p(f"Objectif : {l.get('objectif')}", body_style))
            story.append(_p(l.get("contenu") or "—", body_style))

    story.append(PageBreak())

    story.append(_p("Plans pédagogiques", title_style))
    story.append(Spacer(1, 4 * mm))

    plans = form.get("plans") or []
    if not plans:
        story.append(_p("Aucun plan pédagogique n’est encore rattaché à cette formation.", body_style))
    else:
        for p in plans:
            story.append(_p(f"{p.get('codification') or '—'} • {p.get('titre') or 'Plan pédagogique'}", section_style))
            story.append(_p(f"Modalité générale : {p.get('modalite_generale') or '—'}", body_style))
            story.append(_p(f"Durée cumulée : {p.get('duree_totale') or 0} h • {p.get('nb_blocs') or 0} bloc(s)", body_style))

            for b in p.get("blocs") or []:
                story.append(_p(f"• {b.get('titre') or 'Bloc'} — {b.get('duree') or '—'} h", body_style))

    return story

def _build_plan_pdf_story(form: dict, plan: dict) -> list:
    styles = build_pdf_styles()

    title_style = ParagraphStyle(
        "LearnPlanTitle",
        parent=styles["title"],
        fontName="Helvetica-Bold",
        fontSize=16,
        leading=19,
        textColor=colors.HexColor("#1f2937"),
        spaceAfter=4,
    )

    section_style = ParagraphStyle(
        "LearnPlanSection",
        parent=styles["section"],
        fontName="Helvetica-Bold",
        fontSize=11.2,
        leading=13,
        textColor=colors.HexColor("#c2410c"),
        spaceAfter=5,
        spaceBefore=8,
    )

    body_style = ParagraphStyle(
        "LearnPlanBody",
        parent=styles["body"],
        fontName="Helvetica",
        fontSize=8.8,
        leading=11.2,
        textColor=colors.HexColor("#1f2937"),
        spaceAfter=2,
    )

    small_style = ParagraphStyle(
        "LearnPlanSmall",
        parent=styles["small"],
        fontName="Helvetica",
        fontSize=7.8,
        leading=9.4,
        textColor=colors.HexColor("#6b7280"),
    )

    story = []

    story.append(_p("Plan pédagogique", title_style))
    story.append(_p(f"{form.get('code') or '—'} • {form.get('titre') or 'Formation'}", styles["subtitle"]))
    story.append(Spacer(1, 4 * mm))

    meta = [
        [_p("Référence plan", small_style), _p(plan.get("codification") or "—", body_style), _p("Titre", small_style), _p(plan.get("titre") or "—", body_style)],
        [_p("Modalité générale", small_style), _p(plan.get("modalite_generale") or "—", body_style), _p("Durée cumulée", small_style), _p(f"{plan.get('duree_totale') or 0} h", body_style)],
        [_p("Nombre de blocs", small_style), _p(str(plan.get("nb_blocs") or 0), body_style), _p("Commentaire", small_style), _p(plan.get("commentaire") or "—", body_style)],
    ]

    tbl = Table(meta, colWidths=[30 * mm, 60 * mm, 30 * mm, 60 * mm])
    tbl.setStyle(TableStyle([
        ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor("#e5e7eb")),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e5e7eb")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(tbl)

    story.append(_p("Blocs pédagogiques", section_style))

    blocs = plan.get("blocs") or []
    if not blocs:
        story.append(_p("Aucun bloc pédagogique n’est rattaché à ce plan.", body_style))
    else:
        rows = [[_p("Bloc", small_style), _p("Objectif", small_style), _p("Durée / modalité", small_style)]]

        for b in blocs:
            rows.append([
                _p(b.get("titre") or "—", body_style),
                _p(b.get("objectif") or "—", body_style),
                _p(f"{b.get('duree') or '—'} h • {b.get('modalite_intervention') or '—'}", body_style),
            ])

        table = Table(rows, colWidths=[54 * mm, 82 * mm, 44 * mm])
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#fff7ed")),
            ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#e5e7eb")),
            ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e5e7eb")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        story.append(table)

    return story

@router.get("/learn/formations/{id_effectif}/{id_form}/plans/{id_plan_peda}")
def learn_formation_plan_detail(
    id_effectif: str,
    id_form: str,
    id_plan_peda: str,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        pid = (id_plan_peda or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not pid:
            raise HTTPException(status_code=400, detail="id_plan_peda manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                oid = (profile.get("id_owner") or "").strip()

                return _fetch_plan_detail(cur, oid, fid, pid)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations plan detail error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}/plans")
def learn_formation_plan_create(
    id_effectif: str,
    id_form: str,
    payload: FormationPlanPayload,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        titre = (payload.titre or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not titre:
            raise HTTPException(status_code=400, detail="Titre du plan obligatoire.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                if not _formation_exists_owner(cur, oid, fid):
                    raise HTTPException(status_code=404, detail="Formation introuvable.")

                pid = str(uuid.uuid4())
                code = _next_plan_code(cur, oid)

                cur.execute(
                    """
                    INSERT INTO public.tbl_plan_pedagogique
                      (
                        id_plan_peda,
                        id_owner,
                        id_form,
                        codification,
                        titre,
                        commentaire,
                        modalite_generale,
                        archive,
                        date_creation,
                        date_modification
                      )
                    VALUES
                      (
                        %s, %s, %s,
                        %s, %s, %s, %s,
                        FALSE,
                        NOW(),
                        NOW()
                      )
                    """,
                    (
                        pid,
                        oid,
                        fid,
                        code,
                        titre,
                        _clean_text(payload.commentaire),
                        _clean_text(payload.modalite_generale),
                    ),
                )

                _insert_plan_blocs(cur, oid, fid, pid, payload.blocs)

                item = _fetch_plan_detail(cur, oid, fid, pid)

                conn.commit()

        return {"ok": True, "id_plan_peda": pid, "codification": code, "item": item}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations plan create error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}/plans/{id_plan_peda}")
def learn_formation_plan_update(
    id_effectif: str,
    id_form: str,
    id_plan_peda: str,
    payload: FormationPlanPayload,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        pid = (id_plan_peda or "").strip()
        titre = (payload.titre or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not pid:
            raise HTTPException(status_code=400, detail="id_plan_peda manquant.")
        if not titre:
            raise HTTPException(status_code=400, detail="Titre du plan obligatoire.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                cur.execute(
                    """
                    UPDATE public.tbl_plan_pedagogique
                    SET titre = %s,
                        commentaire = %s,
                        modalite_generale = %s,
                        date_modification = NOW()
                    WHERE id_owner = %s
                      AND id_form = %s
                      AND id_plan_peda = %s
                      AND COALESCE(archive, FALSE) = FALSE
                    """,
                    (
                        titre,
                        _clean_text(payload.commentaire),
                        _clean_text(payload.modalite_generale),
                        oid,
                        fid,
                        pid,
                    ),
                )

                if cur.rowcount <= 0:
                    raise HTTPException(status_code=404, detail="Plan pédagogique introuvable.")

                _archive_plan_blocs_and_sequences(cur, oid, fid, pid)
                _insert_plan_blocs(cur, oid, fid, pid, payload.blocs)

                item = _fetch_plan_detail(cur, oid, fid, pid)

                conn.commit()

        return {"ok": True, "id_plan_peda": pid, "item": item}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations plan update error: {e}")

@router.get("/learn/formations/{id_effectif}/{id_form}/plans/{id_plan_peda}/fiche_pdf")
def learn_formation_plan_fiche_pdf(
    id_effectif: str,
    id_form: str,
    id_plan_peda: str,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        pid = (id_plan_peda or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not pid:
            raise HTTPException(status_code=400, detail="id_plan_peda manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                oid = (profile.get("id_owner") or "").strip()

                form = _fetch_form_detail(cur, oid, fid)
                plan = None

                for p in form.get("plans") or []:
                    if str(p.get("id_plan_peda") or "").strip() == pid:
                        plan = p
                        break

                if not plan:
                    raise HTTPException(status_code=404, detail="Plan pédagogique introuvable.")

                logo_bytes = _fetch_owner_logo_bytes(cur, oid)

        code_label = plan.get("codification") or "Plan"
        titre_label = plan.get("titre") or "Plan pédagogique"
        owner_label = (profile.get("nom_owner") or "Novoskill Learn").strip() or "Novoskill Learn"

        filename = _pdf_latin1_safe(
            f"Plan pédagogique {_pdf_safe_filename_part(code_label, 32)} - {_pdf_safe_filename_part(titre_label, 80)}.pdf"
        )

        pdf_bytes = build_pdf_document(
            _build_plan_pdf_story(form, plan),
            meta={
                "title": _pdf_latin1_safe(f"Plan pédagogique - {code_label} - {titre_label}"),
                "doc_label": _pdf_latin1_safe("Plan pédagogique"),
                "footer_left": _pdf_latin1_safe("Novoskill Learn • Plan pédagogique"),
                "header_right": _pdf_latin1_safe(owner_label),
                "header_right_font_name": "Helvetica-Bold",
                "header_right_font_size": 10.5,
                "logo_bytes": logo_bytes,
            },
        )

        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'inline; filename="{filename}"',
                "Cache-Control": "no-store",
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations plan fiche_pdf error: {e}")


@router.post("/learn/formations/{id_effectif}/{id_form}/plans/{id_plan_peda}/archive")
def learn_formation_plan_archive(
    id_effectif: str,
    id_form: str,
    id_plan_peda: str,
    request: Request,
):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        pid = (id_plan_peda or "").strip()

        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")
        if not pid:
            raise HTTPException(status_code=400, detail="id_plan_peda manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                _learn_require_min_role(profile, "supervisor")
                oid = (profile.get("id_owner") or "").strip()

                cur.execute(
                    """
                    UPDATE public.tbl_plan_pedagogique
                    SET archive = TRUE,
                        date_modification = NOW()
                    WHERE id_owner = %s
                      AND id_form = %s
                      AND id_plan_peda = %s
                      AND COALESCE(archive, FALSE) = FALSE
                    """,
                    (oid, fid, pid),
                )

                if cur.rowcount <= 0:
                    raise HTTPException(status_code=404, detail="Plan pédagogique introuvable.")

                conn.commit()

        return {"ok": True}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations plan archive error: {e}")

@router.get("/learn/formations/{id_effectif}/{id_form}/fiche_pdf")
def learn_formation_fiche_pdf(id_effectif: str, id_form: str, request: Request):
    auth = request.headers.get("Authorization", "")
    u = learn_require_user(auth)

    try:
        fid = (id_form or "").strip()
        if not fid:
            raise HTTPException(status_code=400, detail="id_form manquant.")

        with get_conn() as conn:
            with conn.cursor(row_factory=dict_row) as cur:
                profile = _learn_require_profile(cur, u, id_effectif)
                oid = (profile.get("id_owner") or "").strip()

                form = _fetch_form_detail(cur, oid, fid)
                logo_bytes = _fetch_owner_logo_bytes(cur, oid)

        code_label = form.get("code") or "Formation"
        titre_label = form.get("titre") or "Formation"
        owner_label = (profile.get("nom_owner") or "Novoskill Learn").strip() or "Novoskill Learn"

        filename = _pdf_latin1_safe(
            f"Fiche formation {_pdf_safe_filename_part(code_label, 32)} - {_pdf_safe_filename_part(titre_label, 80)}.pdf"
        )

        pdf_bytes = build_pdf_document(
            _build_formation_pdf_story(form),
            meta={
                "title": _pdf_latin1_safe(f"Fiche formation - {code_label} - {titre_label}"),
                "doc_label": _pdf_latin1_safe("Fiche formation"),
                "footer_left": _pdf_latin1_safe("Novoskill Learn • Fiche formation"),
                "header_right": _pdf_latin1_safe(owner_label),
                "header_right_font_name": "Helvetica-Bold",
                "header_right_font_size": 10.5,
                "logo_bytes": logo_bytes,
            },
        )

        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'inline; filename="{filename}"',
                "Cache-Control": "no-store",
            },
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"learn/formations fiche_pdf error: {e}")